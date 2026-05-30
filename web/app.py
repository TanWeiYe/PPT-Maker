from __future__ import annotations

import json
import os
import re
import shutil
import subprocess
import sys
import uuid
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable
from urllib.parse import urlparse

from flask import Flask, jsonify, render_template, request, send_file
from pptx import Presentation
from pptx.util import Inches, Pt

ROOT_DIR = Path(__file__).resolve().parents[1]
WEB_DIR = ROOT_DIR / "web"
UPLOAD_DIR = WEB_DIR / "uploads"
GENERATED_DIR = WEB_DIR / "generated"
SKILL_DIR = ROOT_DIR / ".agents" / "skills" / "ppt-master"
SCRIPTS_DIR = SKILL_DIR / "scripts"
SAMPLE_DATA_PATH = ROOT_DIR / "examples" / "sample-data.json"

ALLOWED_EXTENSIONS = {
    ".pdf",
    ".doc",
    ".docx",
    ".txt",
    ".md",
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
    ".pptx",
}
SESSION_ID_PATTERN = re.compile(r"^[a-f0-9-]{8,64}$")
MAX_CONTEXT_CHARS = 4000
MAX_OUTLINE_POINT_LENGTH = 36
SUPPORTED_SOURCE_SCRIPTS = {
    "pdf": (SCRIPTS_DIR / "source_to_md" / "pdf_to_md.py").resolve(),
    "doc": (SCRIPTS_DIR / "source_to_md" / "doc_to_md.py").resolve(),
    "web": (SCRIPTS_DIR / "source_to_md" / "web_to_md.py").resolve(),
}


@dataclass
class SessionData:
    session_id: str
    root: Path
    upload_dir: Path
    generated_dir: Path


app = Flask(__name__)


def ensure_dir(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


def normalize_session_id(session_id: str | None = None) -> str:
    if session_id and SESSION_ID_PATTERN.fullmatch(session_id):
        return session_id
    return str(uuid.uuid4())


def init_session(session_id: str | None = None) -> SessionData:
    sid = normalize_session_id(session_id)
    root = WEB_DIR / "runtime" / sid
    upload_dir = root / "uploads"
    generated_dir = root / "generated"
    ensure_dir(upload_dir)
    ensure_dir(generated_dir)
    return SessionData(session_id=sid, root=root, upload_dir=upload_dir, generated_dir=generated_dir)


def run_ppt_master_script(script: Path, script_args: list[str], cwd: Path | None = None) -> tuple[bool, str]:
    if script.resolve() not in SUPPORTED_SOURCE_SCRIPTS.values():
        return False, f"unsupported script path: {script}"
    try:
        result = subprocess.run(
            [sys.executable, str(script), *script_args],
            cwd=str(cwd or ROOT_DIR),
            check=False,
            capture_output=True,
            text=True,
        )
    except OSError as exc:
        return False, str(exc)
    output = "\n".join([result.stdout.strip(), result.stderr.strip()]).strip()
    return result.returncode == 0, output


def is_safe_http_url(url: str) -> bool:
    parsed = urlparse(url)
    return parsed.scheme in {"http", "https"} and bool(parsed.netloc)


def normalize_lines(text: str) -> list[str]:
    return [line.strip() for line in text.splitlines() if line.strip()]


def simple_outline_from_prompt(prompt: str) -> str:
    prompt = prompt.strip()
    if not prompt:
        return "1. 封面\n2. 背景与问题\n3. 核心内容\n4. 总结"

    sentences = re.split(r"[。！？!?.]\s*", prompt)
    points = [s.strip() for s in sentences if s.strip()]
    if not points:
        points = [prompt]

    top = [f"{index}. {point[:MAX_OUTLINE_POINT_LENGTH]}" for index, point in enumerate(points[:8], start=1)]
    if not top:
        top = ["1. 封面", "2. 核心内容", "3. 总结"]
    if not any("总结" in item for item in top):
        top.append(f"{len(top)+1}. 总结")
    return "\n".join(top)


def parse_outline_to_slides(outline_text: str) -> list[tuple[str, list[str]]]:
    lines = normalize_lines(outline_text)
    slides: list[tuple[str, list[str]]] = []
    current_title = ""
    bullets: list[str] = []

    def flush() -> None:
        nonlocal current_title, bullets
        if current_title:
            slides.append((current_title, bullets[:]))
            current_title = ""
            bullets = []

    for line in lines:
        if re.match(r"^(#|\d+[.)]|[一二三四五六七八九十]+[、.])", line):
            flush()
            current_title = re.sub(r"^(#\s*|\d+[.)]\s*|[一二三四五六七八九十]+[、.]\s*)", "", line).strip() or "未命名页面"
        elif line.startswith(("-", "*", "•")):
            bullets.append(line.lstrip("-*• ").strip())
        else:
            if current_title:
                bullets.append(line)
            else:
                current_title = line

    flush()
    if not slides:
        slides = [("封面", ["学生向美观演示"]), ("内容", ["请补充详细大纲"])]
    return slides


def apply_template_if_provided(session: SessionData, template_file: Path | None) -> Path | None:
    if not template_file:
        return None
    target = session.root / "template.pptx"
    shutil.copy2(template_file, target)
    return target


def build_pptx(slides: Iterable[tuple[str, list[str]]], output_path: Path, template_path: Path | None = None) -> None:
    prs = Presentation(str(template_path)) if template_path else Presentation()
    if len(prs.slide_layouts) == 0:
        raise ValueError("PPT 模板中未找到可用版式")
    body_layout_index = 1 if len(prs.slide_layouts) > 1 else 0

    for index, (title, bullets) in enumerate(slides):
        layout = prs.slide_layouts[0] if index == 0 else prs.slide_layouts[body_layout_index]
        slide = prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = title

        if len(slide.placeholders) > 1:
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()
            if bullets:
                text_frame.text = bullets[0]
                text_frame.paragraphs[0].font.size = Pt(24)
                for bullet in bullets[1:]:
                    paragraph = text_frame.add_paragraph()
                    paragraph.text = bullet
                    paragraph.level = 0
                    paragraph.font.size = Pt(20)
            else:
                text_frame.text = ""
        else:
            textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3.5))
            tf = textbox.text_frame
            for bullet in bullets or [""]:
                p = tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(20)

    prs.save(str(output_path))


def save_uploaded_files(session: SessionData, files) -> list[Path]:
    saved: list[Path] = []
    for incoming in files:
        if not incoming.filename:
            continue
        suffix = Path(incoming.filename).suffix.lower()
        if suffix not in ALLOWED_EXTENSIONS:
            continue
        safe_name = f"{uuid.uuid4().hex}{suffix}"
        target = session.upload_dir / safe_name
        incoming.save(target)
        saved.append(target)
    return saved


def detect_template_name(uploaded: list[Path]) -> str | None:
    for path in uploaded:
        if path.suffix.lower() == ".pptx":
            return path.name
    return None


def convert_sources_with_ppt_master(uploaded: list[Path], session: SessionData, source_url: str | None) -> list[Path]:
    converted: list[Path] = []
    for file_path in uploaded:
        suffix = file_path.suffix.lower()
        converter = None
        if suffix == ".pdf":
            converter = SUPPORTED_SOURCE_SCRIPTS["pdf"]
        elif suffix in {".doc", ".docx"}:
            converter = SUPPORTED_SOURCE_SCRIPTS["doc"]

        if converter and converter.exists():
            ok, _ = run_ppt_master_script(converter, [str(file_path)])
            candidate = file_path.with_suffix(".md")
            if ok and candidate.exists():
                converted.append(candidate)
                continue
        if suffix in {".md", ".txt"}:
            converted.append(file_path)

    if source_url and is_safe_http_url(source_url):
        web_converter = SUPPORTED_SOURCE_SCRIPTS["web"]
        if web_converter.exists():
            web_md = session.upload_dir / f"web_{uuid.uuid4().hex}.md"
            ok, _ = run_ppt_master_script(web_converter, [source_url, "-o", str(web_md)], cwd=session.upload_dir)
            if ok and web_md.exists():
                converted.append(web_md)

    return converted


def load_sample_payload() -> dict[str, str]:
    fallback = {
        "title": "人工智能课堂汇报（示例）",
        "prompt": "为高一学生制作一份人工智能入门课堂汇报，包含概念、应用、风险与学习建议。",
        "source_url": "https://en.wikipedia.org/wiki/Artificial_intelligence",
        "outline": "1. 封面：人工智能入门\n2. 人工智能是什么\n3. 典型应用场景\n4. 机遇与风险\n5. 学习建议\n6. 总结",
    }
    if not SAMPLE_DATA_PATH.exists():
        return fallback
    try:
        loaded = json.loads(SAMPLE_DATA_PATH.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        app.logger.warning("sample data unavailable, fallback applied")
        return fallback
    if not isinstance(loaded, dict):
        return fallback
    result: dict[str, str] = {}
    for key, default in fallback.items():
        value = loaded.get(key, default)
        result[key] = value if isinstance(value, str) and value.strip() else default
    return result


@app.get("/")
def index():
    return render_template("index.html")


@app.get("/api/sample")
def sample_payload():
    return jsonify(load_sample_payload())


@app.post("/api/outline")
def generate_outline():
    session_id = request.form.get("session_id")
    if session_id and not SESSION_ID_PATTERN.fullmatch(session_id):
        return jsonify({"error": "invalid session id"}), 400
    prompt = request.form.get("prompt", "")
    source_url = request.form.get("source_url", "").strip() or None

    session = init_session(session_id)
    uploaded = save_uploaded_files(session, request.files.getlist("files"))
    converted = convert_sources_with_ppt_master(uploaded, session, source_url)

    merged_context = [prompt.strip()]
    for file_path in converted:
        try:
            # Keep context bounded to keep request latency stable for web-side MVP interaction.
            content = file_path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            content = file_path.read_text(encoding="utf-8", errors="ignore")
        except OSError:
            continue
        merged_context.append(content[:MAX_CONTEXT_CHARS])

    outline = simple_outline_from_prompt("\n".join([item for item in merged_context if item]))
    payload = {
        "session_id": session.session_id,
        "outline": outline,
        "uploaded_files": [path.name for path in uploaded],
        "template_name": detect_template_name(uploaded),
    }
    (session.root / "outline.json").write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return jsonify(payload)


@app.post("/api/generate")
def generate_pptx():
    data = request.get_json(silent=True) or {}
    session_id = data.get("session_id")
    if session_id and not SESSION_ID_PATTERN.fullmatch(session_id):
        return jsonify({"error": "invalid session id"}), 400
    session = init_session(session_id)

    outline = data.get("outline", "")
    if not outline:
        return jsonify({"error": "outline is required"}), 400

    template_name = data.get("template_name", "")
    template_file = None
    if template_name:
        candidate_template = session.upload_dir / template_name
        if candidate_template.exists():
            template_file = candidate_template
    template_path = apply_template_if_provided(session, template_file)

    slides = parse_outline_to_slides(outline)
    output_path = session.generated_dir / "presentation.pptx"
    try:
        build_pptx(slides, output_path, template_path)
    except OSError:
        app.logger.error("failed to build pptx: OSError")
        return jsonify({"error": "failed to build pptx: file system error"}), 500
    except ValueError:
        app.logger.error("failed to build pptx: ValueError")
        return jsonify({"error": "failed to build pptx: invalid template"}), 500

    return jsonify(
        {
            "session_id": session.session_id,
            "file": output_path.name,
            "download_url": f"/api/download/{session.session_id}/{output_path.name}",
            "quality": {
                "visual": "统一主题和布局基础可用，可继续优化品牌样式",
                "content": "按大纲拆页生成，可继续补充细节",
                "generation": "输出为可编辑 PPTX，支持下载",
            },
        }
    )


@app.get("/api/download/<session_id>/<filename>")
def download(session_id: str, filename: str):
    if not SESSION_ID_PATTERN.fullmatch(session_id):
        return jsonify({"error": "invalid session id"}), 400
    if Path(filename).name != filename or not filename.endswith(".pptx"):
        return jsonify({"error": "invalid filename"}), 400

    base_dir = (WEB_DIR / "runtime" / session_id / "generated").resolve()
    file_path = (base_dir / filename).resolve()
    if not file_path.is_relative_to(base_dir):
        return jsonify({"error": "invalid path"}), 400
    if not file_path.exists():
        return jsonify({"error": "file not found"}), 404
    return send_file(file_path, as_attachment=True)


if __name__ == "__main__":
    ensure_dir(WEB_DIR / "runtime")
    app.run(host="0.0.0.0", port=int(os.environ.get("PORT", "5000")), debug=False)
